import json
import os
import subprocess
from typing import Literal, Optional

import pandas as pd
from mcp.server.fastmcp import FastMCP
from pptx import Presentation

# ==============================================================================
# 🚨 여기가 핵심: 무조건 이 경로 기준으로만 동작하게 강제 설정
# ==============================================================================
BASE_DIR = "/workspaces/clara-ssot"


def get_safe_path(filename: str) -> str:
    """
    무조건 BASE_DIR 아래로 고정시킵니다. 파일명만 입력해도 알아서 경로를 찾아줍니다.
    """
    # 1. 파일명만 남기고 경로 떼어내기 
    clean_name = os.path.basename(filename)

    # 2. 진짜 경로 합치기
    full_path = os.path.join(BASE_DIR, clean_name)

    # (선택) 하위 폴더(Baseline 등)를 지정하고 싶을 때를 위한 예외 처리
    # 입력이 "Baseline/test.xlsx" 형태라면 그대로 유지
    if "/" in filename and not filename.startswith("/"):
        full_path = os.path.join(BASE_DIR, filename)

    # 폴더가 없으면 만들기 (예: Baseline 폴더가 없으면 생성)
    dir_name = os.path.dirname(full_path)
    if dir_name:
        os.makedirs(dir_name, exist_ok=True)

    return full_path


# 서버 인스턴스 생성
mcp = FastMCP("clara-ssot-integrator")

# ==============================================================================
# 1. PPTX Reader
# ==============================================================================


@mcp.tool()
def read_working_deck(filename: str) -> str:
    """
    PPTX 파일을 읽습니다. 파일명만 입력하면 됩니다. (예: 'my_deck.pptx')
    """
    try:
        # 경로 자동 보정
        target_path = get_safe_path(filename)

        # PPTX는 하위 폴더(WorkingDeck/CLARA-SSoT)에 있을 수 있으므로
        # 파일이 없으면 재귀적으로 찾기 (편의성 기능)
        if not os.path.exists(target_path):
            found = False
            for root, dirs, files in os.walk(BASE_DIR):
                if filename in files:
                    target_path = os.path.join(root, filename)
                    found = True
                    break
            if not found:
                return f"Error: {filename} 파일을 {BASE_DIR} 및 하위 폴더에서 찾을 수 없습니다."

        prs = Presentation(target_path)
        full_text = []

        for i, slide in enumerate(prs.slides):
            slide_text = []
            if slide.shapes.title:
                slide_text.append(f"[Title] {slide.shapes.title.text}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape != slide.shapes.title:
                    if shape.text.strip():
                        slide_text.append(shape.text)
            if slide_text:
                full_text.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))

        return "\n\n".join(full_text) if full_text else "내용 없음"

    except Exception as e:
        return f"PPTX 읽기 실패: {str(e)}"


# ==============================================================================
# 2. Excel Sync (자동 경로 보정 적용)
# ==============================================================================


@mcp.tool()
def sync_schema_excel(
    operation: Literal["export_to_excel", "import_from_excel"],
    json_path: str,
    excel_path: str,
) -> str:
    """
    JSON <-> Excel 동기화. 파일명만 넣으면 자동으로 프로젝트 루트에 저장됩니다.
    """
    try:
        # 경로 강제 보정
        real_json_path = get_safe_path(json_path)
        real_excel_path = get_safe_path(excel_path)

        if operation == "export_to_excel":
            # (데이터 생성 로직이 없으면 빈 껍데기라도 만듦 - 테스트용)
            if not os.path.exists(real_json_path):
                # 파일이 없으면 빈 리스트로 시작 (테스트 편의성)
                data = []
            else:
                with open(real_json_path, "r", encoding="utf-8") as f:
                    data = json.load(f)

            if isinstance(data, list):
                df = pd.json_normalize(data)
            else:
                df = pd.DataFrame([data])  # 리스트가 아니면 리스트로 감쌈

            df.to_excel(real_excel_path, index=False)
            return f"✅ 저장 완료! 위치: {real_excel_path} (이제 윈도우 탐색기에서 보입니다)"

        elif operation == "import_from_excel":
            if not os.path.exists(real_excel_path):
                return f"Error: 엑셀 파일({os.path.basename(real_excel_path)})이 없습니다."

            df = pd.read_excel(real_excel_path).fillna("")
            records = df.to_dict(orient="records")

            with open(real_json_path, "w", encoding="utf-8") as f:
                json.dump(records, f, indent=2, ensure_ascii=False)

            return f"✅ 업데이트 완료! 위치: {real_json_path}"

    except Exception as e:
        return f"작업 실패: {str(e)}"


# ==============================================================================
# 3. Git (경로 지정 불필요, 이미 루트에서 실행됨)
# ==============================================================================


@mcp.tool()
def run_git_command(command: str, message: Optional[str] = None) -> str:
    try:
        cmd_list = ["git"] + command.split()
        if "commit" in command and message:
            cmd_list.extend(["-m", message])

        # cwd 옵션으로 무조건 프로젝트 루트에서 실행하게 함
        result = subprocess.run(
            cmd_list, capture_output=True, text=True, check=False, cwd=BASE_DIR
        )

        if result.returncode == 0:
            return f"✅ Success:\n{result.stdout}"
        else:
            return f"❌ Error:\n{result.stderr}"

    except Exception as e:
        return f"Git 에러: {str(e)}"


if __name__ == "__main__":
    mcp.run()
